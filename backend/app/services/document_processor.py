"""
Document Processor Service
Extracts text and images from various document formats
"""

import os
from typing import Tuple, List
from PyPDF2 import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from PIL import Image
from pdf2image import convert_from_path


class DocumentProcessor:
    """Process various document formats"""
    
    def __init__(self):
        self.processors = {
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.xlsx': self._process_xlsx,
            '.pptx': self._process_pptx,
            '.txt': self._process_txt,
            '.jpg': self._process_image,
            '.jpeg': self._process_image,
            '.png': self._process_image
        }
    
    def process_document(self, file_path: str) -> Tuple[str, List[str]]:
        """
        Process document and extract text and images
        Returns: (text, list_of_image_paths)
        """
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext not in self.processors:
            raise ValueError(f"Unsupported format: {ext}")
        
        return self.processors[ext](file_path)
    
    def _process_pdf(self, file_path: str) -> Tuple[str, List[str]]:
        """Extract from PDF"""
        text = ""
        images = []
        
        try:
            # Extract text
            reader = PdfReader(file_path)
            for page in reader.pages:
                text += page.extract_text() + "\n\n"
            
            # Convert pages to images for OCR
            temp_dir = os.path.dirname(file_path)
            page_images = convert_from_path(file_path, dpi=200)
            
            for i, img in enumerate(page_images):
                img_path = os.path.join(temp_dir, f"page_{i+1}.png")
                img.save(img_path, "PNG")
                images.append(img_path)
        except Exception as e:
            print(f"PDF Error: {e}")
        
        return text.strip(), images
    
    def _process_docx(self, file_path: str) -> Tuple[str, List[str]]:
        """Extract from DOCX"""
        text = ""
        images = []
        
        try:
            doc = Document(file_path)
            
            # Extract text
            for para in doc.paragraphs:
                text += para.text + "\n"
            
            # Extract images
            temp_dir = os.path.dirname(file_path)
            for i, rel in enumerate(doc.part.rels.values()):
                if "image" in rel.target_ref:
                    img_data = rel.target_part.blob
                    img_path = os.path.join(temp_dir, f"img_{i+1}.png")
                    with open(img_path, "wb") as f:
                        f.write(img_data)
                    images.append(img_path)
        except Exception as e:
            print(f"DOCX Error: {e}")
        
        return text.strip(), images
    
    def _process_xlsx(self, file_path: str) -> Tuple[str, List[str]]:
        """Extract from Excel"""
        text = ""
        
        try:
            wb = load_workbook(file_path, data_only=True)
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text += f"\n\n=== {sheet_name} ===\n\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(c) if c else "" for c in row])
                    if row_text.strip():
                        text += row_text + "\n"
        except Exception as e:
            print(f"XLSX Error: {e}")
        
        return text.strip(), []
    
    def _process_pptx(self, file_path: str) -> Tuple[str, List[str]]:
        """Extract from PowerPoint"""
        text = ""
        images = []
        
        try:
            prs = Presentation(file_path)
            temp_dir = os.path.dirname(file_path)
            
            for i, slide in enumerate(prs.slides, 1):
                text += f"\n\n=== Slide {i} ===\n\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                    
                    if shape.shape_type == 13:  # Picture
                        try:
                            img_path = os.path.join(temp_dir, f"slide_{i}.png")
                            with open(img_path, "wb") as f:
                                f.write(shape.image.blob)
                            images.append(img_path)
                        except:
                            pass
        except Exception as e:
            print(f"PPTX Error: {e}")
        
        return text.strip(), images
    
    def _process_txt(self, file_path: str) -> Tuple[str, List[str]]:
        """Extract from text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read().strip(), []
        except Exception as e:
            print(f"TXT Error: {e}")
            return "", []
    
    def _process_image(self, file_path: str) -> Tuple[str, List[str]]:
        """Process image file"""
        return "", [file_path]
